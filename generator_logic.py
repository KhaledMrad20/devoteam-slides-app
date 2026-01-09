import json
import re
import os
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Inches, Pt
# On importe load_dotenv seulement si on est en local
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass # Sur Streamlit Cloud, dotenv n'est pas strictement requis si on utilise les Secrets

# --- CONFIGURATION ---
# On récupère la clé depuis les variables d'environnement (Secrets ou .env)
API_KEY = os.getenv("GEMINI_API_KEY")

MODEL_NAME = "gemini-1.5-flash-002"

# --- LAYOUT MAP ---
LAYOUT_MAP = { 
    "COVER": 0,      # Modifié à 0 (souvent Title Slide)
    "SOMMAIRE": 1,
    "SECTION": 2,
    "CONTENT": 1,    # Modifié à 1 (souvent Title and Content)
    "THANK_YOU": 2   # Modifié à 2 (Section Header)
}

# --- TEXT CLEANING ---
def clean_text(text):
    if not text: return "Untitled"
    text = re.sub(r'\s*\((SECTION|CONTENT|Section|Content)\)\s*', '', text, flags=re.IGNORECASE)
    text = re.sub(r'^\d+[\.\-\)\s]+\s*', '', text) 
    text = text.replace('**', '').replace('__', '')
    return text.strip()

def clean_json_response(text):
    text = text.strip()
    match = re.search(r"```json\s*(.*?)\s*```", text, re.DOTALL)
    if match: return match.group(1)
    start = text.find('{')
    end = text.rfind('}')
    if start != -1 and end != -1: return text[start:end+1]
    return text

def extract_strict_sommaire(full_text):
    lines = full_text.split('\n')
    for line in lines:
        match = re.search(r'^(?:Sommaire|Plan)\s*[:\-]\s*(.*)', line, re.IGNORECASE)
        if match:
            raw_items = re.split(r'[,;]', match.group(1))
            return [clean_text(item) for item in raw_items if item.strip()]
    return []

def get_sorted_text_boxes(slide, slide_height):
    shapes = []
    try:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.top > (slide_height * 0.9): continue 
                shapes.append(shape)
        shapes.sort(key=lambda s: s.top)
    except: pass
    return shapes

def safe_add_slide(prs, layout_index):
    try:
        if layout_index < len(prs.slide_layouts):
            return prs.slides.add_slide(prs.slide_layouts[layout_index])
        else:
            return prs.slides.add_slide(prs.slide_layouts[0])
    except:
        return prs.slides.add_slide(prs.slide_layouts[0])

def generate_presentation_outline(content):
    print(f"--- 1. CALLING GEMINI AI ({MODEL_NAME}) ---")
    
    # Check if API Key exists
    if not API_KEY:
        return {"presentation_title": "Error Occurred", "subtitle": "Clé API manquante dans les Secrets Streamlit.", "slides": []}

    is_short_topic = len(content.strip()) < 150
    
    json_structure = """
    {
        "presentation_title": "The Title Here",
        "subtitle": "The Subtitle Here",
        "slides": [
            { "title": "Section Name", "content": ["Bullet 1", "Bullet 2"] }
        ]
    }
    """
    
    if is_short_topic:
        prompt = f"""
        You are a Presentation Creator. TOPIC: "{content}"
        INSTRUCTIONS: Create professional slides. Output valid JSON only.
        STRUCTURE: {json_structure}
        """
    else:
        prompt = f"""
        You are a Data Extractor. INPUT TEXT: "{content}"
        INSTRUCTIONS: Extract content exactly. Output valid JSON only.
        STRUCTURE: {json_structure}
        """

    try:
        client = genai.Client(api_key=API_KEY)
        response = client.models.generate_content(
            model=MODEL_NAME,
            contents=prompt,
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )
        
        if not response.text: raise ValueError("Réponse vide de l'IA")

        raw_text = clean_json_response(response.text)
        data = json.loads(raw_text)
        data['original_text'] = content
        data['mode'] = 'creative' if is_short_topic else 'strict'
        return data

    except Exception as e:
        print(f"⚠️ GEMINI ERROR: {str(e)}")
        return { "presentation_title": "Error Occurred", "subtitle": str(e), "slides": [] }

def create_presentation_file(data, template_path="my_brand_template.pptx", output_filename="presentation_generee.pptx"):
    print(f"--- 2. GENERATING PPTX ---")
    try: prs = Presentation(template_path)
    except: 
        print("Template not found, using blank.")
        prs = Presentation()
    
    h = prs.slide_height
    w = prs.slide_width
    
    main_title = data.get('presentation_title', 'Untitled Presentation')
    subtitle = data.get('subtitle', '')
    original_text = data.get('original_text', '')
    mode = data.get('mode', 'strict')
    
    # --- STEP 1: PREPARE PLAN ---
    final_summary_items = []
    if mode == 'strict':
        final_summary_items = extract_strict_sommaire(original_text)
    
    if not final_summary_items:
        seen = set()
        for s in data.get('slides', []):
            t = clean_text(s.get('title', ''))
            if t and t not in seen:
                final_summary_items.append(t)
                seen.add(t)

    # --- STEP 2: GROUP SLIDES ---
    grouped_slides = {}
    ordered_topics = []
    
    for slide_item in data.get('slides', []):
        raw_title = clean_text(slide_item.get('title', 'Untitled'))
        key = raw_title.lower()
        if key not in grouped_slides:
            grouped_slides[key] = { "title": raw_title, "content": [] }
            ordered_topics.append(key)
        
        raw_content = slide_item.get('content', [])
        if isinstance(raw_content, str): raw_content = [raw_content]
        elif isinstance(raw_content, list): pass
        else: raw_content = []
        grouped_slides[key]["content"].extend(raw_content)

    # --- STEP 3: BUILD SLIDES ---
    # 1. COVER
    try:
        slide = safe_add_slide(prs, LAYOUT_MAP["COVER"])
        boxes = get_sorted_text_boxes(slide, h)
        if len(boxes) > 0: boxes[0].text_frame.text = clean_text(main_title)
        if len(boxes) > 1: boxes[1].text_frame.text = subtitle
    except: pass

    # 2. PLAN
    try:
        slide = safe_add_slide(prs, LAYOUT_MAP["SOMMAIRE"])
        valid_boxes = get_sorted_text_boxes(slide, h)
        if len(valid_boxes) > 0: valid_boxes[0].text_frame.text = "PLAN"
        
        list_box = None
        for box in valid_boxes[1:]: 
            if (box.left + (box.width / 2)) < (w / 2):
                list_box = box
                break
        if list_box is None and len(valid_boxes) > 1: list_box = valid_boxes[1]

        if list_box:
            list_box.text_frame.clear()
            if list_box.width < Inches(4): list_box.width = Inches(4.5)
            summary_text = []
            for i, item in enumerate(final_summary_items, 1):
                summary_text.append(f"{i}. {item}")
            list_box.text_frame.text = "\n".join(summary_text)
    except: pass

    # 3. SECTIONS & CONTENT
    section_count = 0
    for topic_key in ordered_topics:
        topic_data = grouped_slides[topic_key]
        final_title = topic_data['title']
        final_content = topic_data['content']
        section_count += 1
        
        # Content Slide
        try:
            slide = safe_add_slide(prs, LAYOUT_MAP["CONTENT"])
            boxes = get_sorted_text_boxes(slide, h)
            if len(boxes) > 0: boxes[0].text_frame.text = final_title
            
            tf = None
            if len(boxes) > 1:
                tf = boxes[1].text_frame
            else:
                left = Inches(1); top = Inches(2); width = Inches(8); height = Inches(4)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame; tf.word_wrap = True

            if tf:
                tf.clear()
                for item in final_content:
                    p = tf.add_paragraph()
                    p.text = str(item)
                    p.level = 0
                    p.font.size = Pt(18)
                    p.space_after = Pt(10)
        except: pass

    # 4. THANK YOU
    try:
        safe_add_slide(prs, LAYOUT_MAP["THANK_YOU"])
    except: pass

    prs.save(output_filename)
    return output_filename



